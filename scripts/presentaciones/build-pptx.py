"""Arma el .pptx de una presentación a partir de sus láminas ya renderizadas.

Cada lámina entra como imagen a sangre completa y las notas del presentador salen
del atributo data-notes del HTML, así que en PowerPoint las notas sí son texto
editable aunque el diseño de la lámina sea una imagen.

Uso:
    python scripts/presentaciones/build-pptx.py <index.html> <dir-imagenes> <salida.pptx>

El procedimiento completo (cómo generar las imágenes) está en
docs/presentacion-ia-2026.md. Requiere `pip install python-pptx`.
"""

import glob
import html
import os
import re
import sys

from pptx import Presentation
from pptx.util import Emu, Inches


def attr(blob: str, name: str) -> str:
    m = re.search(rf'{name}="([^"]*)"', blob)
    return html.unescape(m.group(1)) if m else ""


def main() -> None:
    if len(sys.argv) != 4:
        raise SystemExit(__doc__)
    html_path, img_dir, out_path = sys.argv[1:4]

    source = open(html_path, encoding="utf-8").read()
    sections = re.findall(r'<section class="slide[^"]*"([^>]*)>', source)
    meta = [(attr(b, "data-title"), attr(b, "data-notes")) for b in sections]

    images = sorted(glob.glob(os.path.join(img_dir, "slide-*.jpg")))
    if not images:
        raise SystemExit(f"no hay slide-*.jpg en {img_dir}")
    if len(images) != len(meta):
        raise SystemExit(
            f"{len(images)} imágenes vs {len(meta)} láminas en el HTML: no coinciden"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for path, (title, notes) in zip(images, meta):
        slide = prs.slides.add_slide(blank)
        pic = slide.shapes.add_picture(path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
        if title:
            # Nombra la imagen con el título: es lo que se ve en el panel de selección.
            pic.name = title[:120]
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.core_properties.title = "IA 2026 · De conversar a dirigir"
    prs.core_properties.author = "Capital Academy"
    prs.core_properties.comments = (
        "Exportado desde https://capitalacademy.cl/presentaciones/ia-2026 · "
        "Cada lámina es una imagen; el texto editable vive en las notas del presentador."
    )
    prs.save(out_path)
    size_mb = os.path.getsize(out_path) / 1_048_576
    print(f"{len(images)} láminas → {out_path} ({size_mb:.1f} MB)")


if __name__ == "__main__":
    main()
